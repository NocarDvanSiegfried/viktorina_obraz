from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.db.models import Question, Quiz

TITLE_COLOR = RGBColor(31, 41, 55)
ACCENT_COLOR = RGBColor(37, 99, 235)
SUCCESS_COLOR = RGBColor(21, 128, 61)
MUTED_COLOR = RGBColor(75, 85, 99)


class QuizPptxExportService:
    def _set_title_style(self, text_frame, *, size: int = 36) -> None:
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = "Calibri"
        run.font.bold = True
        run.font.size = Pt(size)
        run.font.color.rgb = TITLE_COLOR

    def _set_body_style(self, paragraph, *, size: int = 18, color: RGBColor = MUTED_COLOR) -> None:
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color

    def _add_intro_slide(self, prs: Presentation, quiz: Quiz, questions: list[Question]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = quiz.title or "Quiz"
        self._set_title_style(title.text_frame)

        subtitle.text = (
            f"Subject: {quiz.subject or '-'} | Grade: {quiz.grade or '-'}\n"
            f"Difficulty: {quiz.difficulty or '-'} | Questions: {len(questions)}"
        )
        subtitle.text_frame.paragraphs[0].line_spacing = 1.3
        self._set_body_style(subtitle.text_frame.paragraphs[0], size=16)

    def _add_agenda_slide(self, prs: Presentation, questions: list[Question]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "План викторины"
        self._set_title_style(slide.shapes.title.text_frame, size=30)

        body = slide.placeholders[1].text_frame
        body.clear()
        for idx, question in enumerate(questions, start=1):
            paragraph = body.add_paragraph() if idx > 1 else body.paragraphs[0]
            paragraph.text = f"{idx}. {question.question_text}"
            paragraph.level = 0
            self._set_body_style(paragraph, size=18)

    def _add_question_slide(
        self,
        prs: Presentation,
        index: int,
        question: Question,
        *,
        include_answers: bool,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Вопрос {index}"
        self._set_title_style(slide.shapes.title.text_frame, size=30)

        question_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.3))
        question_frame = question_box.text_frame
        question_frame.text = f"[{question.question_type}] {question.question_text}"
        self._set_body_style(question_frame.paragraphs[0], size=20, color=TITLE_COLOR)

        options_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(8.7), Inches(2.6))
        options_frame = options_box.text_frame
        options_frame.clear()
        for option_idx, option in enumerate(question.answers or [], start=1):
            marker = (
                "✓"
                if include_answers and option in (question.correct_answers or [])
                else "•"
            )
            paragraph = (
                options_frame.paragraphs[0]
                if option_idx == 1
                else options_frame.add_paragraph()
            )
            paragraph.text = f"{marker} {option}"
            self._set_body_style(
                paragraph,
                size=18,
                color=SUCCESS_COLOR if marker == "✓" else MUTED_COLOR,
            )

        if include_answers:
            explanation = question.explanation or "Без пояснения."
            explanation_box = slide.shapes.add_textbox(
                Inches(9.75), Inches(3.0), Inches(2.3), Inches(2.8)
            )
            explanation_frame = explanation_box.text_frame
            explanation_frame.text = "Комментарий"
            self._set_body_style(explanation_frame.paragraphs[0], size=14, color=ACCENT_COLOR)
            paragraph = explanation_frame.add_paragraph()
            paragraph.text = explanation
            self._set_body_style(paragraph, size=13)

    def _add_final_slide(self, prs: Presentation) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Готово к уроку"
        self._set_title_style(title.text_frame)
        subtitle.text = "Экспорт PPTX сформирован. Удачного проведения викторины!"
        self._set_body_style(subtitle.text_frame.paragraphs[0], size=18, color=ACCENT_COLOR)

    def build_quiz_pptx(
        self,
        quiz: Quiz,
        questions: list[Question],
        *,
        include_answers: bool = True,
    ) -> bytes:
        prs = Presentation()
        self._add_intro_slide(prs, quiz, questions)
        self._add_agenda_slide(prs, questions)
        for index, question in enumerate(questions, start=1):
            self._add_question_slide(
                prs,
                index,
                question,
                include_answers=include_answers,
            )
        self._add_final_slide(prs)

        output = BytesIO()
        prs.save(output)
        return output.getvalue()


pptx_export_service = QuizPptxExportService()
