from io import BytesIO

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from app.schemas.material import SourceFragment
from app.services.gigachat_service import gigachat_service

PLACEHOLDER_TEXTS = {"string", "source_text", "null", "none"}

EMPTY_PDF_MESSAGE = (
    "Не удалось извлечь текст из PDF. Вставьте текст вручную или загрузите TXT."
)

EMPTY_IMAGE_MESSAGE = (
    "Не удалось распознать текст на изображении. "
    "Вставьте текст вручную или загрузите TXT."
)

SUPPORTED_FILE_EXTENSIONS = (
    ".txt",
    ".pdf",
    ".docx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
)


class MaterialService:
    def _split_text_to_chunks(self, text: str, max_chunk_chars: int = 600) -> list[str]:
        raw_parts = [part.strip() for part in text.split("\n") if part.strip()]
        chunks: list[str] = []
        current_chunk = ""

        for part in raw_parts:
            if len(current_chunk) + len(part) + 1 <= max_chunk_chars:
                current_chunk = f"{current_chunk}\n{part}".strip()
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = part

        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def extract_text_from_txt(self, content: bytes) -> str:
        return content.decode("utf-8", errors="ignore").strip()

    def extract_text_from_pdf(self, content: bytes) -> list[SourceFragment]:
        reader = PdfReader(BytesIO(content))
        fragments: list[SourceFragment] = []

        for page_idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if not text or not text.strip():
                continue

            chunks = self._split_text_to_chunks(text)

            for chunk_idx, chunk_text in enumerate(chunks, start=1):
                fragments.append(
                    SourceFragment(
                        fragment_id=f"pdf_page_{page_idx}_chunk_{chunk_idx}",
                        source_type="pdf",
                        source_name=f"page_{page_idx}_chunk_{chunk_idx}",
                        text=chunk_text,
                    )
                )

        return fragments

    def extract_text_from_docx(self, filename: str, content: bytes) -> list[SourceFragment]:
        document = Document(BytesIO(content))
        fragments: list[SourceFragment] = []

        collected_lines: list[str] = []
        for paragraph in document.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                collected_lines.append(text)

        full_text = "\n".join(collected_lines).strip()
        if not full_text:
            return []

        chunks = self._split_text_to_chunks(full_text)
        for chunk_idx, chunk_text in enumerate(chunks, start=1):
            fragments.append(
                SourceFragment(
                    fragment_id=f"docx_chunk_{chunk_idx}",
                    source_type="docx",
                    source_name=filename,
                    text=chunk_text,
                )
            )

        return fragments

    def extract_text_from_pptx(self, filename: str, content: bytes) -> list[SourceFragment]:
        presentation = Presentation(BytesIO(content))
        fragments: list[SourceFragment] = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())

            slide_text = "\n".join(lines).strip()
            if not slide_text:
                continue

            chunks = self._split_text_to_chunks(slide_text)
            for chunk_idx, chunk_text in enumerate(chunks, start=1):
                fragments.append(
                    SourceFragment(
                        fragment_id=f"pptx_slide_{slide_idx}_chunk_{chunk_idx}",
                        source_type="pptx",
                        source_name=filename,
                        text=chunk_text,
                    )
                )

        return fragments

    def _fragments_from_plain_text(
        self,
        text: str,
        source_type: str,
        filename: str,
        fragment_id_prefix: str,
    ) -> list[SourceFragment]:
        cleaned = text.strip()
        if not cleaned:
            return []

        chunks = self._split_text_to_chunks(cleaned)
        fragments: list[SourceFragment] = []
        for chunk_idx, chunk_text in enumerate(chunks, start=1):
            fragments.append(
                SourceFragment(
                    fragment_id=f"{fragment_id_prefix}_{chunk_idx}",
                    source_type=source_type,
                    source_name=filename,
                    text=chunk_text,
                )
            )
        return fragments

    def extract_fragments_via_vision(
        self,
        filename: str,
        content: bytes,
        source_type: str,
        fragment_id_prefix: str,
    ) -> list[SourceFragment]:
        extracted = gigachat_service.extract_text_from_visual(filename, content)
        return self._fragments_from_plain_text(
            extracted,
            source_type,
            filename,
            fragment_id_prefix,
        )

    def extract_fragments(self, filename: str, content: bytes) -> tuple[str, list[SourceFragment]]:
        lower_name = filename.lower()

        if lower_name.endswith(".txt"):
            text = self.extract_text_from_txt(content)
            if not text:
                return "txt", []

            return "txt", [
                SourceFragment(
                    fragment_id="txt_1",
                    source_type="txt",
                    source_name=filename,
                    text=text,
                )
            ]

        if lower_name.endswith(".pdf"):
            fragments = self.extract_text_from_pdf(content)
            if not fragments:
                fragments = self.extract_fragments_via_vision(
                    filename,
                    content,
                    source_type="pdf",
                    fragment_id_prefix="pdf_ocr",
                )
            return "pdf", fragments

        if lower_name.endswith(".docx"):
            return "docx", self.extract_text_from_docx(filename, content)

        if lower_name.endswith(".pptx"):
            return "pptx", self.extract_text_from_pptx(filename, content)

        if lower_name.endswith((".png", ".jpg", ".jpeg", ".webp")):
            fragments = self.extract_fragments_via_vision(
                filename,
                content,
                source_type="image",
                fragment_id_prefix="image_ocr",
            )
            return "image", fragments

        raise ValueError(
            "Неподдерживаемый формат. Доступны: ручной текст, "
            ".txt, .pdf, .docx, .pptx, .png, .jpg, .jpeg, .webp"
        )

    def merge_fragments(
        self,
        manual_text: str | None,
        file_fragments: list[SourceFragment],
    ) -> list[SourceFragment]:
        fragments: list[SourceFragment] = []

        cleaned_manual = (manual_text or "").strip()
        if cleaned_manual.lower() in PLACEHOLDER_TEXTS:
            cleaned_manual = ""

        if cleaned_manual:
            fragments.append(
                SourceFragment(
                    fragment_id="manual_1",
                    source_type="manual_text",
                    source_name="teacher_input",
                    text=cleaned_manual,
                )
            )

        fragments.extend(file_fragments)
        return fragments

    def build_combined_context(
        self,
        fragments: list[SourceFragment],
        max_chars: int = 6000,
    ) -> str:
        if not fragments:
            return ""

        parts = []
        for fragment in fragments:
            parts.append(
                f"[fragment_id={fragment.fragment_id}; "
                f"source_type={fragment.source_type}; "
                f"source_name={fragment.source_name}]\n{fragment.text}"
            )

        return "\n\n".join(parts).strip()[:max_chars]


material_service = MaterialService()
