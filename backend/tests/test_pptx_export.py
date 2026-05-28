"""Quiz export to PPTX."""

from io import BytesIO
from unittest.mock import patch

from pptx import Presentation

SAMPLE_MODEL_JSON = """
[JSON_START]
{
  "quiz_title": "PPTX Export Quiz",
  "subject": "History",
  "grade": "9",
  "topic": "World War",
  "questions": [
    {
      "type": "single_choice",
      "text": "When did WW2 end?",
      "options": ["1943", "1945", "1947", "1950"],
      "correct_answers": ["1945"],
      "explanation": "War ended in 1945.",
      "difficulty": "easy",
      "source_fragment_id": "manual_1"
    }
  ]
}
[JSON_END]
"""


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_export_pptx_returns_binary(_mock_chat, client):
    owner_id = "owner-day21-pptx"
    created = client.post(
        "/quiz/generate-from-materials",
        data={
            "owner_id": owner_id,
            "subject": "History",
            "grade": "9",
            "topic": "World War",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
            "source_text": "WW2 ended in 1945.",
        },
    )
    assert created.status_code == 200
    quiz_id = created.json()["quiz_id"]

    resp = client.get(f"/quiz/{quiz_id}/export-pptx", params={"owner_id": owner_id})
    assert resp.status_code == 200
    assert (
        resp.headers["content-type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert resp.content.startswith(b"PK")


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_export_pptx_other_owner_403(_mock_chat, client):
    owner_id = "owner-day21-pptx-a"
    other_owner = "owner-day21-pptx-b"
    created = client.post(
        "/quiz/generate-from-materials",
        data={
            "owner_id": owner_id,
            "subject": "History",
            "grade": "9",
            "topic": "World War",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
            "source_text": "WW2 ended in 1945.",
        },
    )
    quiz_id = created.json()["quiz_id"]

    resp = client.get(f"/quiz/{quiz_id}/export-pptx", params={"owner_id": other_owner})
    assert resp.status_code == 403


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_export_pptx_contains_title_and_question_slide(_mock_chat, client):
    owner_id = "owner-day21-pptx-structure"
    created = client.post(
        "/quiz/generate-from-materials",
        data={
            "owner_id": owner_id,
            "subject": "History",
            "grade": "9",
            "topic": "World War",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
            "source_text": "WW2 ended in 1945.",
        },
    )
    assert created.status_code == 200
    quiz_id = created.json()["quiz_id"]

    resp = client.get(f"/quiz/{quiz_id}/export-pptx", params={"owner_id": owner_id})
    assert resp.status_code == 200

    presentation = Presentation(BytesIO(resp.content))
    assert len(presentation.slides) >= 3

    first_slide_texts = [
        shape.text.strip()
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text") and shape.text.strip()
    ]
    assert any("PPTX Export Quiz" in text for text in first_slide_texts)

    question_slide_texts = [
        shape.text.strip()
        for shape in presentation.slides[2].shapes
        if hasattr(shape, "text") and shape.text.strip()
    ]
    joined = "\n".join(question_slide_texts)
    assert "When did WW2 end?" in joined
    assert "1945" in joined


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_export_pptx_classroom_mode_hides_answers(_mock_chat, client):
    owner_id = "owner-day21-pptx-classroom"
    created = client.post(
        "/quiz/generate-from-materials",
        data={
            "owner_id": owner_id,
            "subject": "History",
            "grade": "9",
            "topic": "World War",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
            "source_text": "WW2 ended in 1945.",
        },
    )
    assert created.status_code == 200
    quiz_id = created.json()["quiz_id"]

    resp = client.get(
        f"/quiz/{quiz_id}/export-pptx",
        params={"owner_id": owner_id, "include_answers": "false"},
    )
    assert resp.status_code == 200

    presentation = Presentation(BytesIO(resp.content))
    question_slide_texts = [
        shape.text.strip()
        for shape in presentation.slides[2].shapes
        if hasattr(shape, "text") and shape.text.strip()
    ]
    joined = "\n".join(question_slide_texts)
    assert "When did WW2 end?" in joined
    assert "✓" not in joined
    assert "Комментарий" not in joined
