"""Day 10: DOCX/PPTX parsing and generation."""

from io import BytesIO
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from app.services.material_service import material_service

SAMPLE_MODEL_JSON = """
[JSON_START]
{
  "quiz_title": "Викторина из Office",
  "subject": "Биология",
  "grade": "8",
  "topic": "Клетка",
  "questions": [
    {
      "type": "single_choice",
      "text": "Что является базовой единицей жизни?",
      "options": ["Клетка", "Орган", "Ткань", "Атом"],
      "correct_answers": ["Клетка"],
      "explanation": "Клетка — базовая структурная единица.",
      "difficulty": "easy",
      "source_fragment_id": "docx_chunk_1"
    }
  ]
}
[JSON_END]
"""


def _make_docx_bytes(text: str) -> bytes:
    document = Document()
    document.add_paragraph(text)
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _make_pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Тема"
    slide.placeholders[1].text = text
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_fragments_from_docx_file():
    docx_bytes = _make_docx_bytes("Клетка — базовая единица организма.")
    file_type, fragments = material_service.extract_fragments("lesson.docx", docx_bytes)
    assert file_type == "docx"
    assert len(fragments) >= 1
    assert fragments[0].source_type == "docx"


def test_extract_fragments_from_pptx_file():
    pptx_bytes = _make_pptx_bytes("Клетка имеет мембрану и ядро.")
    file_type, fragments = material_service.extract_fragments("lesson.pptx", pptx_bytes)
    assert file_type == "pptx"
    assert len(fragments) >= 1
    assert fragments[0].source_type == "pptx"


MINIMAL_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01"
    b"\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
)

OCR_TEXT = "Клетка — основная единица жизни. Мембрана защищает содержимое клетки."


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
@patch(
    "app.services.gigachat_service.gigachat_service.extract_text_from_visual",
    return_value=OCR_TEXT,
)
def test_generate_from_png_image(_mock_vision, _mock_chat, client):
    response = client.post(
        "/quiz/generate-from-materials",
        files={"file": ("image.png", BytesIO(MINIMAL_PNG), "image/png")},
        data={
            "owner_id": "owner-day13-image-office",
            "subject": "Биология",
            "grade": "8",
            "topic": "Клетка",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
        },
    )
    assert response.status_code == 200
    assert response.json()["quiz_id"]


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_generate_from_docx_file(_mock_chat, client, db_engine):
    from sqlalchemy.orm import sessionmaker

    from app.db.models import MaterialUpload

    docx_bytes = _make_docx_bytes("Клетка — базовая единица организма.")
    response = client.post(
        "/quiz/generate-from-materials",
        files={
            "file": (
                "lesson.docx",
                BytesIO(docx_bytes),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
        data={
            "owner_id": "owner-day10-docx",
            "subject": "Биология",
            "grade": "8",
            "topic": "Клетка",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
        },
    )
    assert response.status_code == 200
    quiz_id = response.json()["quiz_id"]

    SessionLocal = sessionmaker(bind=db_engine)
    session = SessionLocal()
    try:
        uploads = (
            session.query(MaterialUpload)
            .filter(MaterialUpload.quiz_id == quiz_id)
            .all()
        )
        assert len(uploads) == 1
        assert uploads[0].source_type == "docx"
    finally:
        session.close()


@patch("app.services.gigachat_service.gigachat_service.chat", return_value=SAMPLE_MODEL_JSON)
def test_generate_from_pptx_file(_mock_chat, client, db_engine):
    from sqlalchemy.orm import sessionmaker

    from app.db.models import MaterialUpload

    pptx_bytes = _make_pptx_bytes("Клетка имеет мембрану и ядро.")
    response = client.post(
        "/quiz/generate-from-materials",
        files={
            "file": (
                "lesson.pptx",
                BytesIO(pptx_bytes),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
        data={
            "owner_id": "owner-day10-pptx",
            "subject": "Биология",
            "grade": "8",
            "topic": "Клетка",
            "question_count": "1",
            "difficulty": "easy",
            "question_types": "single_choice",
        },
    )
    assert response.status_code == 200
    quiz_id = response.json()["quiz_id"]

    SessionLocal = sessionmaker(bind=db_engine)
    session = SessionLocal()
    try:
        uploads = (
            session.query(MaterialUpload)
            .filter(MaterialUpload.quiz_id == quiz_id)
            .all()
        )
        assert len(uploads) == 1
        assert uploads[0].source_type == "pptx"
    finally:
        session.close()

